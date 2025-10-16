"""
PPTX Content Parser - Extract content with font information from Standard Template PPTX
"""

import os
from dataclasses import dataclass
from typing import List, Optional

from pptx import Presentation as PythonPPTX
from pptx.enum.shapes import MSO_SHAPE_TYPE

from pptagent.utils import get_logger

logger = get_logger(__name__)


@dataclass
class FontInfo:
    """Font information for text"""
    name: str
    size: float  # in points
    bold: bool = False
    italic: bool = False
    underline: bool = False
    color_rgb: Optional[tuple] = None  # (R, G, B)


@dataclass
class TextRun:
    """A text run with its font information"""
    text: str
    font: FontInfo


@dataclass
class TextBlock:
    """A block of text (paragraph) with multiple runs"""
    runs: List[TextRun]
    level: int = 0  # Indentation level
    
    @property
    def full_text(self) -> str:
        """Get the full text of this block"""
        return "".join(run.text for run in self.runs)
    
    @property
    def dominant_font(self) -> FontInfo:
        """Get the most common font in this block"""
        if not self.runs:
            return FontInfo(name="Arial", size=18)
        # Return the font of the longest run
        longest_run = max(self.runs, key=lambda r: len(r.text))
        return longest_run.font


@dataclass
class SlideContent:
    """Content extracted from a single slide"""
    slide_number: int
    title: Optional[str] = None
    title_font: Optional[FontInfo] = None
    text_blocks: List[TextBlock] = None
    images: List[str] = None  # Image paths if extracted
    
    def __post_init__(self):
        if self.text_blocks is None:
            self.text_blocks = []
        if self.images is None:
            self.images = []


@dataclass
class PPTXContent:
    """Complete content from Standard Template PPTX"""
    slides: List[SlideContent]
    source_file: str
    
    @property
    def num_slides(self) -> int:
        return len(self.slides)
    
    def to_markdown(self) -> str:
        """Convert to markdown format for compatibility"""
        md_lines = []
        
        for slide in self.slides:
            md_lines.append(f"\n## Slide {slide.slide_number}")
            
            if slide.title:
                md_lines.append(f"\n### {slide.title}")
                if slide.title_font:
                    md_lines.append(
                        f"<!-- Font: {slide.title_font.name}, "
                        f"Size: {slide.title_font.size}pt -->"
                    )
            
            for block in slide.text_blocks:
                if block.full_text.strip():
                    indent = "  " * block.level
                    md_lines.append(f"{indent}- {block.full_text}")
                    font = block.dominant_font
                    md_lines.append(
                        f"{indent}<!-- Font: {font.name}, "
                        f"Size: {font.size}pt -->"
                    )
            
            for img_path in slide.images:
                md_lines.append(f"\n![Image]({img_path})")
        
        return "\n".join(md_lines)


def extract_font_info(run) -> FontInfo:
    """Extract font information from a text run"""
    font = run.font
    
    # Get font name
    font_name = font.name if font.name else "Arial"
    
    # Get font size (in points)
    font_size = font.size.pt if font.size else 18
    
    # Get font style
    bold = font.bold if font.bold is not None else False
    italic = font.italic if font.italic is not None else False
    underline = font.underline if font.underline is not None else False
    
    # Get color
    color_rgb = None
    try:
        if font.color and font.color.rgb:
            color_rgb = (
                font.color.rgb[0],
                font.color.rgb[1],
                font.color.rgb[2]
            )
    except:
        pass
    
    return FontInfo(
        name=font_name,
        size=font_size,
        bold=bold,
        italic=italic,
        underline=underline,
        color_rgb=color_rgb
    )


def extract_text_from_shape(shape) -> List[TextBlock]:
    """Extract text blocks from a shape"""
    text_blocks = []
    
    if not hasattr(shape, "text_frame"):
        return text_blocks
    
    text_frame = shape.text_frame
    
    for paragraph in text_frame.paragraphs:
        runs = []
        
        for run in paragraph.runs:
            if run.text.strip():  # Only process non-empty runs
                font_info = extract_font_info(run)
                runs.append(TextRun(text=run.text, font=font_info))
        
        if runs:  # Only add blocks with content
            level = paragraph.level if hasattr(paragraph, 'level') else 0
            text_blocks.append(TextBlock(runs=runs, level=level))
    
    return text_blocks


async def parse_standard_template_pptx(pptx_path: str, output_folder: str) -> PPTXContent:
    """
    Parse a Standard Template PPTX file and extract content with font information.
    
    Args:
        pptx_path: Path to the Standard Template PPTX file
        output_folder: Folder to save extracted images (if any)
    
    Returns:
        PPTXContent: Parsed content with font information
    """
    logger.info(f"Parsing Standard Template PPTX: {pptx_path}")
    
    # Create output folder
    os.makedirs(output_folder, exist_ok=True)
    images_folder = os.path.join(output_folder, "images")
    os.makedirs(images_folder, exist_ok=True)
    
    # Load presentation
    prs = PythonPPTX(pptx_path)
    
    slides_content = []
    
    for slide_idx, slide in enumerate(prs.slides, start=1):
        logger.info(f"Processing slide {slide_idx}/{len(prs.slides)}")
        
        slide_content = SlideContent(slide_number=slide_idx)
        
        for shape in slide.shapes:
            # Extract title
            if shape.is_placeholder:
                placeholder_format = shape.placeholder_format
                if placeholder_format.type == 1:  # Title placeholder
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.title = shape.text
                        # Get title font
                        if hasattr(shape, "text_frame") and shape.text_frame.paragraphs:
                            first_para = shape.text_frame.paragraphs[0]
                            if first_para.runs:
                                slide_content.title_font = extract_font_info(first_para.runs[0])
                    continue
            
            # Extract text content
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX or \
               shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER or \
               hasattr(shape, "text_frame"):
                text_blocks = extract_text_from_shape(shape)
                slide_content.text_blocks.extend(text_blocks)
            
            # Extract images
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    image_bytes = image.blob
                    ext = image.ext
                    image_filename = f"slide_{slide_idx}_image_{len(slide_content.images)}.{ext}"
                    image_path = os.path.join(images_folder, image_filename)
                    
                    with open(image_path, 'wb') as f:
                        f.write(image_bytes)
                    
                    slide_content.images.append(image_path)
                    logger.info(f"Extracted image: {image_filename}")
                except Exception as e:
                    logger.warning(f"Failed to extract image from slide {slide_idx}: {e}")
        
        slides_content.append(slide_content)
    
    logger.info(f"Successfully parsed {len(slides_content)} slides from Standard Template")
    
    return PPTXContent(slides=slides_content, source_file=pptx_path)


def content_to_document_format(content: PPTXContent) -> dict:
    """
    Convert PPTXContent to Document format for compatibility with existing pipeline.
    
    Args:
        content: Parsed PPTX content
    
    Returns:
        dict: Document-like structure
    """
    sections = []
    
    for slide in content.slides:
        section = {
            "heading": slide.title or f"Slide {slide.slide_number}",
            "content": [],
            "metadata": {
                "slide_number": slide.slide_number,
                "title_font": {
                    "name": slide.title_font.name if slide.title_font else "Arial",
                    "size": slide.title_font.size if slide.title_font else 24
                }
            }
        }
        
        # Add text blocks
        for block in slide.text_blocks:
            font = block.dominant_font
            section["content"].append({
                "type": "text",
                "text": block.full_text,
                "font": {
                    "name": font.name,
                    "size": font.size,
                    "bold": font.bold,
                    "italic": font.italic,
                    "underline": font.underline
                },
                "level": block.level
            })
        
        # Add images
        for img_path in slide.images:
            section["content"].append({
                "type": "image",
                "path": img_path
            })
        
        sections.append(section)
    
    return {
        "sections": sections,
        "source_type": "pptx",
        "num_slides": content.num_slides
    }

